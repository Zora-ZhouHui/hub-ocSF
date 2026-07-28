"""
Harness 核心引擎 - 渐进式披露架构

基于 Harness Engineering 设计理念：
- 常驻层: Skill 索引始终加载
- 触发层: 匹配条件后按需加载完整 Skill
- 执行层: Skill 执行期间完整驻留
"""

import os
import re
import time
from pathlib import Path
from typing import Dict, List, Optional


class ProgressiveDisclosure:
    """渐进式披露管理器"""

    def __init__(self, config_dir: str):
        self.config_dir = Path(config_dir)
        self.skill_index = self._load_skill_index()
        self.loaded_skills: Dict[str, dict] = {}

    def _load_skill_index(self) -> List[dict]:
        """加载常驻层: Skill 索引"""
        index_path = self.config_dir / "skill_index.md"
        if not index_path.exists():
            return []

        skills = []
        content = index_path.read_text(encoding="utf-8")
        current_skill = None

        for line in content.split("\n"):
            line = line.strip()
            if line.startswith("## ") and line.strip() == "## 可用 Skills":
                continue
            if line.startswith("- `"):
                if current_skill:
                    skills.append(current_skill)
                current_skill = {"triggers": [], "description": ""}
                match = re.match(r"- `([^`]+)`\s*—\s*(.*)", line)
                if match:
                    current_skill["name"] = match.group(1)
                    current_skill["description"] = match.group(2)
            elif current_skill and "触发词:" in line:
                triggers_str = line.split("触发词:")[1].strip()
                current_skill["triggers"] = [t.strip() for t in triggers_str.split("|")]
            elif current_skill and "描述:" in line:
                current_skill["description"] = line.split("描述:")[1].strip()

        if current_skill:
            skills.append(current_skill)

        return skills

    def match_skill(self, user_input: str) -> Optional[dict]:
        """根据用户输入匹配 Skill"""
        user_input_lower = user_input.lower()
        for skill in self.skill_index:
            for trigger in skill.get("triggers", []):
                if trigger.lower() in user_input_lower:
                    return skill
        return None

    def load_skill(self, skill_name: str) -> Optional[str]:
        """按需加载完整 Skill 定义"""
        if skill_name in self.loaded_skills:
            return self.loaded_skills[skill_name].get("full_content", "")

        skill_path = self.config_dir.parent / "skills" / skill_name / "SKILL.md"
        if not skill_path.exists():
            return None

        content = skill_path.read_text(encoding="utf-8")
        self.loaded_skills[skill_name] = {
            "full_content": content,
            "loaded_at": time.time(),
        }
        return content

    def release_skill(self, skill_name: str):
        """释放已加载的 Skill"""
        self.loaded_skills.pop(skill_name, None)

    def get_index_summary(self) -> str:
        """获取常驻层摘要"""
        lines = ["# Skill 索引（常驻层）\n"]
        for skill in self.skill_index:
            triggers = ", ".join(skill.get("triggers", [])[:3])
            lines.append(f"- `{skill['name']}` — {skill.get('description', '')}")
            lines.append(f"  触发词: {triggers}")
        return "\n".join(lines)


class MemorySystem:
    """四层记忆模型"""

    def __init__(self, config_dir: str):
        self.config_dir = Path(config_dir)
        self.working_memory: List[dict] = []
        self.long_term_memory = self._load_long_term_memory()

    def _load_long_term_memory(self) -> dict:
        """加载 L3 长期记忆"""
        memory = {}
        mem_path = self.config_dir / "MEMORY.md"
        if mem_path.exists():
            content = mem_path.read_text(encoding="utf-8")
            memory["raw"] = content
            memory["sections"] = self._parse_memory_sections(content)
        return memory

    def _parse_memory_sections(self, content: str) -> Dict[str, str]:
        """解析 MEMORY.md 的各个段落"""
        sections = {}
        current_section = None
        for line in content.split("\n"):
            if line.startswith("## "):
                current_section = line[3:].strip()
                sections[current_section] = []
            elif current_section and line.strip():
                sections[current_section].append(line)
        return sections

    def add_to_working(self, role: str, content: str):
        """添加到 L1 工作记忆"""
        self.working_memory.append(
            {"role": role, "content": content, "timestamp": time.time()}
        )

    def get_context(self) -> str:
        """组装上下文（四层检索策略）"""
        parts = []

        if self.long_term_memory.get("sections"):
            for section, lines in self.long_term_memory["sections"].items():
                if section in ["Skill 索引（常驻层）", "用户偏好", "技术栈"]:
                    parts.append(f"## {section}\n" + "\n".join(lines))

        recent = self.working_memory[-10:]
        if recent:
            parts.append("## 最近对话")
            for msg in recent:
                parts.append(f"[{msg['role']}]: {msg['content']}")

        return "\n\n".join(parts)

    def flush(self):
        """记忆刷新"""
        if self.working_memory:
            print(f"[Memory Flush] 刷新 {len(self.working_memory)} 条工作记忆")
            self.working_memory = []


class Gateway:
    """Fat Gateway: 消息网关与会话管理"""

    def __init__(self):
        self.sessions: Dict[str, dict] = {}
        self.lane_queues: Dict[str, List[dict]] = {}

    def create_session(self, session_id: str, user_id: str = "default") -> dict:
        """创建新会话"""
        session = {
            "sessionId": session_id,
            "userId": user_id,
            "channel": "harness",
            "createdAt": time.time(),
            "isRunning": False,
            "hasError": False,
            "retryCount": 0,
        }
        self.sessions[session_id] = session
        self.lane_queues[session_id] = []
        return session

    def submit_message(self, session_id: str, content: str) -> dict:
        """提交消息到 Lane 队列"""
        if session_id not in self.sessions:
            self.create_session(session_id)

        lane = self.lane_queues[session_id]
        message = {"content": content, "timestamp": time.time(), "processed": False}
        lane.append(message)
        print(f"[Gateway] 消息入队 (session={session_id}, queue_size={len(lane)})")
        return message

    def get_next_message(self, session_id: str) -> Optional[dict]:
        """获取队列中下一条消息"""
        if session_id not in self.lane_queues:
            return None

        lane = self.lane_queues[session_id]
        for msg in lane:
            if not msg["processed"]:
                return msg
        return None

    def mark_processed(self, session_id: str, message: dict):
        """标记消息为已处理"""
        message["processed"] = True

    def get_session_status(self, session_id: str) -> dict:
        """获取会话状态"""
        if session_id not in self.sessions:
            return {"error": "Session not found"}
        session = self.sessions[session_id]
        lane = self.lane_queues.get(session_id, [])
        pending = [m for m in lane if not m["processed"]]
        return {
            "sessionId": session_id,
            "isRunning": session["isRunning"],
            "pendingCount": len(pending),
            "retryCount": session["retryCount"],
        }


class AgentNode:
    """Agent 执行引擎: ReAct 循环"""

    def __init__(
        self, gateway: Gateway, memory: MemorySystem, disclosure: ProgressiveDisclosure
    ):
        self.gateway = gateway
        self.memory = memory
        self.disclosure = disclosure
        self.max_turns = 10
        self.tools = self._register_tools()
        self.base_dir = Path(__file__).resolve().parent.parent
        self.result_dir = self.base_dir / "result"
        self.result_dir.mkdir(parents=True, exist_ok=True)

    def _register_tools(self) -> dict:
        """注册可用工具"""
        return {
            "ppt_extractor": self._tool_ppt_extractor,
            "knowledge_summarizer": self._tool_knowledge_summarizer,
            "web_searcher": self._tool_web_searcher,
            "markdown_writer": self._tool_markdown_writer,
        }

    def run(self, session_id: str, user_input: str) -> dict:
        """执行 ReAct 循环"""
        print(f"\n[AgentNode] 开始执行 (session={session_id})")
        print(f"[AgentNode] 用户输入: {user_input[:100]}...\n")

        self.memory.add_to_working("user", user_input)

        matched_skill = self.disclosure.match_skill(user_input)
        if not matched_skill:
            return self._handle_no_match(user_input)

        skill_content = self.disclosure.load_skill(matched_skill["name"])
        if not skill_content:
            return {
                "status": "error",
                "message": f"Skill {matched_skill['name']} 加载失败",
            }

        print(f"[Progressive Disclosure] 匹配 Skill: {matched_skill['name']}")
        print(f"[Progressive Disclosure] 按需加载完整定义 ({len(skill_content)} chars)")

        result = self._execute_skill(matched_skill["name"], skill_content, user_input)

        self.disclosure.release_skill(matched_skill["name"])
        print(f"[Progressive Disclosure] 释放 Skill {matched_skill['name']}")

        return result

    def _execute_skill(
        self, skill_name: str, skill_content: str, user_input: str
    ) -> dict:
        """执行 Skill 流程（ReAct 循环），共享 execution_context"""
        steps = self._parse_skill_steps(skill_content)
        results = []
        execution_context = {
            "user_input": user_input,
            "ppt_files": [],
            "extracted_texts": {},
            "summary": {},
            "supplementary_topics": [],
            "output_path": None,
        }

        for i, step in enumerate(steps):
            print(f"\n[ReAct] 第 {i + 1} 步: {step['name']}")
            print(f"[Reason] {step.get('reasoning', '')[:100]}...")

            tool_name = step.get("tool", "")
            if tool_name in self.tools:
                tool_result = self.tools[tool_name](execution_context, step)
                results.append(
                    {"step": step["name"], "tool": tool_name, "result": tool_result}
                )
                print(f"[Act] 调用 {tool_name}")
                print(f"[Observation] {str(tool_result)[:200]}...")
            else:
                results.append(
                    {"step": step["name"], "tool": None, "result": "Skipped"}
                )

        return {
            "status": "success",
            "skill": skill_name,
            "steps_completed": len(results),
            "results": results,
            "output_path": execution_context.get("output_path"),
            "summary": execution_context.get("summary"),
        }

    def _parse_skill_steps(self, skill_content: str) -> List[dict]:
        """解析 Skill 执行步骤"""
        steps = []
        if "ppt-knowledge-extractor" in skill_content:
            steps = [
                {
                    "name": "PPT 文字提取",
                    "reasoning": "需要先从 PPT 文件中提取原始文字内容",
                    "tool": "ppt_extractor",
                },
                {
                    "name": "知识总结",
                    "reasoning": "从提取的内容中提炼核心概念和架构",
                    "tool": "knowledge_summarizer",
                },
                {
                    "name": "补充缺失知识",
                    "reasoning": "搜索初学者可能缺失的前置知识",
                    "tool": "web_searcher",
                },
                {
                    "name": "生成 Markdown",
                    "reasoning": "将总结内容输出为结构化文档",
                    "tool": "markdown_writer",
                },
            ]
        return steps

    def _tool_ppt_extractor(self, ctx: dict, step: dict) -> dict:
        """工具: PPT 文字提取"""
        try:
            from pptx import Presentation
        except ImportError:
            return {
                "status": "error",
                "message": "python-pptx 未安装，请执行: pip install python-pptx",
            }

        ppt_paths = self._extract_ppt_paths(ctx["user_input"])
        if not ppt_paths:
            return {"status": "error", "message": "未找到 PPT 文件路径"}

        ctx["ppt_files"] = ppt_paths
        extracted = {}
        for path in ppt_paths:
            try:
                prs = Presentation(path)
                slides = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    texts.append(text)
                    if texts:
                        slides.append({"page": i + 1, "content": texts})
                extracted[path] = slides
                print(f"  [PPT Extractor] {Path(path).name}: {len(slides)} 页提取成功")
            except Exception as e:
                extracted[path] = {"error": str(e)}

        ctx["extracted_texts"] = extracted
        return {
            "status": "success",
            "files_processed": len(ppt_paths),
            "total_slides": sum(
                len(v) for v in extracted.values() if isinstance(v, list)
            ),
        }

    def _tool_knowledge_summarizer(self, ctx: dict, step: dict) -> dict:
        """工具: 知识总结 - 基于提取内容生成真实总结"""
        extracted = ctx.get("extracted_texts", {})
        if not extracted:
            return {"status": "error", "message": "没有可总结的内容"}

        all_texts = []
        for filepath, slides in extracted.items():
            if isinstance(slides, list):
                for slide in slides:
                    all_texts.extend(slide.get("content", []))

        combined_text = "\n".join(all_texts)

        summary = self._analyze_and_summarize(combined_text)
        ctx["summary"] = summary

        topic = self._infer_topic(all_texts)
        ctx["topic"] = topic

        print(f"  [Knowledge Summarizer] 分析 {len(all_texts)} 条文本，生成总结")
        print(f"  [Knowledge Summarizer] 主题推断: {topic}")

        return {
            "status": "success",
            "core_concepts_found": len(summary.get("core_concepts", [])),
            "topic": topic,
        }

    def _analyze_and_summarize(self, text: str) -> dict:
        """分析文本并生成结构化总结"""
        concepts = []
        architectures = []
        comparisons = []
        best_practices = []
        applications = []

        concept_keywords = [
            "Agent",
            "智能体",
            "ReAct",
            "Function Call",
            "MCP",
            "RAG",
            "Skills",
            "Harness",
            "Prompt Engineering",
            "Context Engineering",
            "记忆模型",
            "Gateway",
            "Fat Gateway",
            "Channel Adapter",
            "Progressive Disclosure",
            "渐进式披露",
            "Context Window",
            "Embedding",
            "向量数据库",
            "Token",
            "Compaction",
            "Memory Flush",
        ]

        for kw in concept_keywords:
            if kw.lower() in text.lower():
                concepts.append(kw)

        if any(
            k in text for k in ["五大核心组件", "四大", "三层", "四层", "架构", "组件"]
        ):
            architectures.append("分层架构模型")
        if any(
            k in text for k in ["Gateway", "Nodes", "Skills", "Memory", "Control UI"]
        ):
            architectures.append("OpenClaw 五大组件")
        if any(k in text for k in ["四层记忆", "L1", "L2", "L3", "L4"]):
            architectures.append("四层记忆模型")
        if any(k in text for k in ["渐进式披露", "Progressive Disclosure"]):
            architectures.append("渐进式披露架构")

        if any(k in text for k in ["对比", "vs", "区别", "不同于", "差异"]):
            comparisons.append("技术对比分析")
        if any(k in text for k in ["Platform vs Framework", "Function Call vs ReAct"]):
            comparisons.append("框架对比")

        if any(k in text for k in ["原则", "最佳实践", "设计原则", "自检"]):
            best_practices.append("设计原则")
        if any(k in text for k in ["单一职责", "触发词精准", "自检机制"]):
            best_practices.append("Skill 设计原则")

        if any(k in text for k in ["应用", "场景", "案例", "个人助手", "企业"]):
            applications.append("应用场景")

        supplementary = []
        concept_set = set(concepts)
        if "Agent" in concept_set or "智能体" in concept_set:
            supplementary.append("AI Agent 基础概念")
        if "MCP" in concept_set:
            supplementary.append("MCP 协议详解")
        if "RAG" in concept_set:
            supplementary.append("RAG 检索增强生成")
        if "向量数据库" in concept_set or "Embedding" in concept_set:
            supplementary.append("向量数据库选型")
        if "ReAct" in concept_set:
            supplementary.append("ReAct 模式详解")
        if "Function Call" in concept_set:
            supplementary.append("Function Call 机制")

        return {
            "core_concepts": concepts,
            "architectures": architectures,
            "comparisons": comparisons,
            "best_practices": best_practices,
            "applications": applications,
            "supplementary_topics": supplementary,
            "raw_text_length": len(text),
        }

    def _infer_topic(self, texts: list) -> str:
        """推断主题"""
        all_text = " ".join(texts)
        if any(k in all_text for k in ["Harness", "OpenClaw", "Agent 平台"]):
            return "AI_Agent_Harness_Architecture"
        if any(k in all_text for k in ["Skills", "Function Call", "MCP", "RAG"]):
            return "AI_Agent_Skills_Paradigm"
        if any(k in all_text for k in ["记忆", "Memory", "记忆模型"]):
            return "AI_Agent_Memory_System"
        return "AI_Agent_Knowledge_Summary"

    def _tool_web_searcher(self, ctx: dict, step: dict) -> dict:
        """工具: 网络搜索补充 - 识别需补充的知识点"""
        summary = ctx.get("summary", {})
        topics = summary.get("supplementary_topics", [])

        search_queries = []
        for topic in topics:
            query = f"{topic} 入门教程 详解 2025"
            search_queries.append(query)

        if not search_queries:
            search_queries = [
                "AI Agent 入门教程 基础概念详解",
                "ReAct Function Call 区别 对比",
                "MCP 协议 架构 详解",
                "RAG 检索增强生成 原理",
            ]

        ctx["supplementary_topics"] = topics
        print(f"  [Web Searcher] 规划 {len(search_queries)} 个搜索查询")
        for q in search_queries[:3]:
            print(f"    - {q}")

        return {
            "status": "success",
            "queries_planned": search_queries,
            "topics_to_supplement": topics,
        }

    def _tool_markdown_writer(self, ctx: dict, step: dict) -> dict:
        """工具: Markdown 文档生成 - 写入 week13/result/"""
        summary = ctx.get("summary", {})
        topic = ctx.get("topic", "AI_Agent_Knowledge_Summary")
        supplementary = ctx.get("supplementary_topics", [])
        ppt_files = ctx.get("ppt_files", [])

        md_content = self._build_markdown(summary, topic, supplementary, ppt_files)

        safe_name = topic.replace(" ", "_").replace("/", "_")
        timestamp = time.strftime("%Y%m%d_%H%M%S")
        filename = f"{safe_name}_{timestamp}.md"

        output_path = self.result_dir / filename
        output_path.write_text(md_content, encoding="utf-8")

        abs_path = os.path.abspath(str(output_path))
        ctx["output_path"] = abs_path

        print(f"  [Markdown Writer] 文件已生成: {abs_path}")
        print(f"  [Markdown Writer] 文件大小: {len(md_content)} chars")

        return {
            "status": "success",
            "output_path": abs_path,
            "file_size": len(md_content),
            "filename": filename,
        }

    def _build_markdown(
        self, summary: dict, topic: str, supplementary: list, ppt_files: list
    ) -> str:
        """构建 Markdown 内容"""
        lines = []
        lines.append(f"# {topic.replace('_', ' ')} 核心知识体系总结\n")
        lines.append(f"> 生成时间: {time.strftime('%Y-%m-%d %H:%M:%S')}")
        lines.append(f"> 来源: {', '.join(Path(f).name for f in ppt_files)}")
        lines.append("")

        lines.append("## 一、基础概念补充\n")
        lines.append("### 1.1 什么是 AI Agent（智能体）\n")
        lines.append(
            "**AI Agent** 是能够**自主感知环境、规划任务、调用工具并执行**的AI系统。"
        )
        lines.append(
            '它不是单一模型，而是一套 **"大模型(LLM) + 记忆 + 规划 + 工具"** 的组合拳。\n'
        )
        lines.append("| 对比维度 | 传统AI（聊天机器人） | AI Agent |")
        lines.append("|---|---|---|")
        lines.append("| 工作模式 | **被动响应** | **主动推进** |")
        lines.append("| 处理任务 | 单一步骤 | 复杂多步 |")
        lines.append("| 记忆能力 | 短期记忆 | 长期记忆 |")
        lines.append("| 遇到问题 | 卡住就停 | 自我修正 |\n")

        lines.append("### 1.2 核心术语\n")
        lines.append("| 术语 | 解释 |")
        lines.append("|---|---|")
        lines.append("| **Token** | LLM 处理文本的最小单位 |")
        lines.append("| **Context Window** | LLM 单次能处理的最大 token 数 |")
        lines.append("| **Embedding** | 将文本转换为高维向量的过程 |")
        lines.append("| **Vector Database** | 专门存储和检索向量的数据库 |")
        lines.append("| **幻觉 (Hallucination)** | LLM 生成不准确或虚构内容的现象 |\n")

        lines.append("## 二、核心内容总结\n")
        concepts = summary.get("core_concepts", [])
        if concepts:
            lines.append(f"### 2.1 核心概念 ({len(concepts)} 个)\n")
            for concept in concepts:
                lines.append(f"- **{concept}**")
            lines.append("")

        archs = summary.get("architectures", [])
        if archs:
            lines.append(f"### 2.2 架构模型 ({len(archs)} 个)\n")
            for arch in archs:
                lines.append(f"- {arch}")
            lines.append("")

        comps = summary.get("comparisons", [])
        if comps:
            lines.append(f"### 2.3 技术对比 ({len(comps)} 项)\n")
            for comp in comps:
                lines.append(f"- {comp}")
            lines.append("")

        practices = summary.get("best_practices", [])
        if practices:
            lines.append(f"### 2.4 最佳实践 ({len(practices)} 项)\n")
            for p in practices:
                lines.append(f"- {p}")
            lines.append("")

        apps = summary.get("applications", [])
        if apps:
            lines.append(f"### 2.5 应用场景 ({len(apps)} 类)\n")
            for app in apps:
                lines.append(f"- {app}")
            lines.append("")

        lines.append("## 三、架构与模型详解\n")
        lines.append("### 3.1 工程三层架构\n")
        lines.append(
            "从 Prompt Engineering → Context Engineering → Harness Engineering，外延逐层扩大：\n"
        )
        lines.append("| 层次 | 核心问题 | 技术手段 | 适用场景 |")
        lines.append("|---|---|---|---|")
        lines.append(
            "| **Prompt Engineering** | 如何把任务说清楚 | Role Prompting, CoT, ReAct | 单轮问答、意图对齐 |"
        )
        lines.append(
            "| **Context Engineering** | 模型关键时刻看到什么 | RAG, Tools, Memory, State | 多工具、长对话、知识密集 |"
        )
        lines.append(
            "| **Harness Engineering** | 模型在什么机制里运行 | 目标边界, 反馈回路, 记录系统 | 生产级、长任务、多Agent |\n"
        )

        lines.append("### 3.2 渐进式披露架构\n")
        lines.append("| 层级 | 内容 | Token 占用 | 加载时机 |")
        lines.append("|---|---|---|---|")
        lines.append("| **常驻层** | Skill 索引摘要 | < 200 | 始终加载 |")
        lines.append("| **触发层** | 完整 Skill 定义 | 500-2000 | 匹配触发词后 |")
        lines.append("| **执行层** | Skill 完整驻留 | 完整 | 执行期间 |\n")

        lines.append("## 四、概念辨析与对比\n")
        lines.append("### 4.1 Function Call vs MCP vs RAG vs Skills\n")
        lines.append("| 概念 | 核心特征 | 类比 |")
        lines.append("|---|---|---|")
        lines.append(
            '| **Function Call** | 模型输出结构化JSON，框架路由执行 | "函数调用" |'
        )
        lines.append('| **MCP** | 标准化工具接口，动态发现与调用 | "USB接口" |')
        lines.append('| **RAG** | 外部知识检索注入，弥补LLM知识缺失 | "外脑" |')
        lines.append('| **Skills** | 行为封装+知识内聚+按需加载 | "专家模块" |\n')
        lines.append(
            "> 四者**不互斥**：Skills 是上层调度者，Function Call、MCP、RAG 是其可选执行模块\n"
        )

        lines.append("### 4.2 四层记忆模型\n")
        lines.append("| 层级 | 名称 | 存储 | 速度 | 生命周期 |")
        lines.append("|---|---|---|---|---|")
        lines.append("| L1 | Working Memory | Context Window | 即时 | 当前会话 |")
        lines.append("| L2 | Short-term Memory | SQLite/每日MD | ~5ms | 当天~数天 |")
        lines.append("| L3 | Long-term Memory | MEMORY.md + USER.md | ~1ms | 永久 |")
        lines.append("| L4 | Vector Search | 向量数据库 | ~10-50ms | 永久 |\n")

        lines.append("## 五、补充知识点\n")
        lines.append("以下为初学者需要掌握的前置知识：\n")
        for i, topic_item in enumerate(supplementary, 1):
            lines.append(f"### 5.{i} {topic_item}\n")
            lines.append("（建议通过搜索引擎进一步学习相关内容）\n")

        vector_db_keywords = ["向量数据库", "Embedding", "Vector Database"]
        has_vector_db_topic = any(
            any(kw in topic for kw in vector_db_keywords) for topic in supplementary
        )
        if not has_vector_db_topic:
            lines.append(f"### 5.{len(supplementary) + 1} 向量数据库选型\n")
            lines.append("| 数据库 | 特点 | 适用场景 |")
            lines.append("|---|---|---|")
            lines.append("| **FAISS** | Facebook开源，高性能 | 学术研究、单机场景 |")
            lines.append("| **Chroma** | 轻量、零依赖 | 开发测试、小型应用 |")
            lines.append("| **Milvus** | 分布式、大规模 | 企业级生产环境 |")
            lines.append("| **Pinecone** | 全托管云服务 | 企业级、免运维 |\n")

        lines.append("## 六、学习路径建议\n")
        lines.append("```")
        lines.append("入门：理解 LLM、Token、Context、幻觉 →")
        lines.append("进阶：学习 Function Call、ReAct、RAG 基础 →")
        lines.append("深入：掌握 MCP 协议、向量数据库、记忆系统 →")
        lines.append("实战：使用 Harness / LangChain 构建 Agent →")
        lines.append("精通：Harness Engineering、Skills 设计、生产级部署")
        lines.append("```\n")

        lines.append("---\n")
        lines.append(f"*文档由 Harness Engine v1.0 自动生成*")
        lines.append(f"*基于 {', '.join(Path(f).name for f in ppt_files)} 的内容*")

        return "\n".join(lines)

    def _extract_ppt_paths(self, user_input: str) -> List[str]:
        """从用户输入中提取 PPT 文件路径"""
        paths = []
        words = user_input.split()
        for word in words:
            clean = word.strip("\"'(),")
            if clean.endswith((".ppt", ".pptx")) and os.path.exists(clean):
                paths.append(clean)
        return paths

    def _handle_no_match(self, user_input: str) -> dict:
        """无匹配 Skill 时的处理"""
        return {
            "status": "no_match",
            "message": "未匹配到可用 Skill",
            "available_skills": [s["name"] for s in self.disclosure.skill_index],
            "hint": "请尝试提供 PPT 文件路径",
        }


class Harness:
    """
    Harness 主引擎

    架构: Gateway → Progressive Disclosure → Agent Node
    记忆: 四层记忆模型 (Working → Short-term → Long-term → Vector)
    披露: 渐进式披露 (常驻层 → 触发层 → 执行层)
    """

    def __init__(self, config_dir: str):
        self.config_dir = config_dir
        self.gateway = Gateway()
        self.memory = MemorySystem(config_dir)
        self.disclosure = ProgressiveDisclosure(config_dir)
        self.agent = AgentNode(self.gateway, self.memory, self.disclosure)
        self._display_banner()

    def _display_banner(self):
        """显示 Harness 启动信息"""
        print("=" * 60)
        print("  Harness Engine v1.0")
        print("  渐进式披露架构 | Progressive Disclosure")
        print("=" * 60)
        print(f"\n  配置目录: {self.config_dir}")
        print(f"  结果输出: {self.agent.result_dir}")
        print(f"  可用 Skills: {len(self.disclosure.skill_index)}")
        for s in self.disclosure.skill_index:
            triggers = ", ".join(s.get("triggers", [])[:2])
            print(f"    • {s['name']} (触发: {triggers})")
        print()

    def process_input(self, user_input: str, session_id: str = None) -> dict:
        """处理用户输入"""
        if not session_id:
            session_id = f"session_{int(time.time())}"

        self.gateway.submit_message(session_id, user_input)
        result = self.agent.run(session_id, user_input)
        self.memory.flush()

        return result

    def get_status(self) -> dict:
        """获取 Harness 状态"""
        return {
            "sessions": len(self.gateway.sessions),
            "skills_loaded": list(self.disclosure.loaded_skills.keys()),
            "working_memory_size": len(self.memory.working_memory),
            "available_skills": [s["name"] for s in self.disclosure.skill_index],
            "result_dir": str(self.agent.result_dir),
        }

    def interactive_mode(self):
        """交互式模式"""
        print("\n进入交互模式，输入 'exit' 退出\n")
        while True:
            try:
                user_input = input("你: ").strip()
                if user_input.lower() in ["exit", "quit", "q", "退出"]:
                    print("\n[Harness] 退出交互模式")
                    break
                if not user_input:
                    continue

                result = self.process_input(user_input)
                self._display_result(result)

            except KeyboardInterrupt:
                print("\n\n[Harness] 中断，退出交互模式")
                break
            except Exception as e:
                print(f"\n[Harness] 错误: {e}")

    def _display_result(self, result: dict):
        """显示执行结果"""
        status = result.get("status", "unknown")
        print(f"\n[Result] 状态: {status}")

        if status == "success":
            print(f"  Skill: {result.get('skill', 'N/A')}")
            print(f"  完成步骤: {result.get('steps_completed', 0)}")
            for r in result.get("results", []):
                tool_result = r.get("result", {})
                if isinstance(tool_result, dict):
                    status_icon = "✓" if tool_result.get("status") == "success" else "✗"
                    extra = ""
                    if tool_result.get("output_path"):
                        extra = f" → {tool_result['output_path']}"
                    elif tool_result.get("file_size"):
                        extra = f" ({tool_result['file_size']} chars)"
                    print(
                        f"    [{status_icon}] {r['step']}: {r.get('tool', 'N/A')}{extra}"
                    )
                else:
                    print(f"    - {r['step']}: {r.get('tool', 'N/A')}")
            if result.get("output_path"):
                print(f"\n  📄 生成文件: {result['output_path']}")
        elif status == "no_match":
            print(f"  可用 Skills: {result.get('available_skills', [])}")
            print(f"  提示: {result.get('hint', '')}")
        elif status == "error":
            print(f"  错误: {result.get('message', '')}")

        print()
